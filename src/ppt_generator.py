"""
PPT 생성기 — 템플릿 슬라이드를 직접 수정하는 방식.

동작 원리:
  1. 템플릿 .pptx 파일을 열어 기존 슬라이드 디자인을 그대로 유지
  2. 각 슬라이드의 텍스트 도형만 분석 결과로 교체
  3. 슬라이드 수 < 섹션 수 → 부족한 슬라이드 복제하여 추가
  4. 템플릿 없으면 기본 다크 테마로 폴백

매핑 규칙 (10장 템플릿 기준):
  슬라이드 1   → 표지  (제목 + 날짜/메타)
  슬라이드 2   → 개요  (키워드 · AI 모델 · 기사 수)
  슬라이드 3~  → 분석 섹션 (10개)
  마지막 슬라이드 → 원본 유지 (감사/클로징)
"""
import copy
import os
import re
import sys
import textwrap
from datetime import datetime, timezone
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def _safe_print(msg: str) -> None:
    """Windows cp949 콘솔에서 특수문자 출력 시 깨짐 방지."""
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode(sys.stdout.encoding or 'cp949', errors='replace').decode(sys.stdout.encoding or 'cp949'))

# ── 기본 다크 테마 색상 ──────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCA, 0xE9, 0xFF)
C_BODY_BG = RGBColor(0x16, 0x21, 0x33)


# ══════════════════════════════════════════════════════════════════════════════
# 템플릿 탐색
# ══════════════════════════════════════════════════════════════════════════════

def _find_template() -> Path | None:
    env_path = os.environ.get("TEMPLATE_PATH", "")
    root = Path(__file__).parent.parent
    candidates = []
    if env_path:
        candidates.append(Path(env_path))
        candidates.append(root / env_path)
    candidates += [
        root / "templates" / "report_template.pptx",
        root / "templates" / "template.pptx",
    ]
    for p in candidates:
        if p.exists():
            return p
    return None


# ══════════════════════════════════════════════════════════════════════════════
# 텍스트 처리 유틸
# ══════════════════════════════════════════════════════════════════════════════

def _clean_md(text: str) -> str:
    """마크다운 기호 제거 → 순수 텍스트."""
    text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", text)
    text = re.sub(r"\*{1,2}([^*]+)\*{1,2}", r"\1", text)
    text = re.sub(r"^#{1,4}\s*", "", text, flags=re.MULTILINE)
    return text.strip()


def _format_body(text: str, max_lines: int = 13) -> str:
    """섹션 본문을 슬라이드용 bullet 텍스트로 변환."""
    lines = []
    for raw in _clean_md(text).split("\n"):
        s = raw.strip().lstrip("-•·").strip()
        if not s:
            continue
        # 길면 줄바꿈 (90자 기준)
        for wrapped in textwrap.wrap(s, width=88):
            lines.append("• " + wrapped)
        if len(lines) >= max_lines:
            break
    return "\n".join(lines[:max_lines])


# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 내 텍스트 도형 탐색 & 교체
# ══════════════════════════════════════════════════════════════════════════════

def _text_shapes_by_position(slide) -> list:
    """텍스트가 있는 도형을 위→아래, 왼→오른 순으로 정렬."""
    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text.strip()
        if txt:
            result.append((shape.top, shape.left, len(txt), shape))
    result.sort(key=lambda x: (x[0], x[1]))
    return [s for *_, s in result]


def _extract_first_rpr(tf) -> etree._Element | None:
    """첫 번째 런의 rPr(글꼴 서식) 요소 반환."""
    try:
        txBody = tf._txBody
        paras  = txBody.findall(qn("a:p"))
        if not paras:
            return None
        runs = paras[0].findall(qn("a:r"))
        if not runs:
            return None
        rpr = runs[0].find(qn("a:rPr"))
        return copy.deepcopy(rpr) if rpr is not None else None
    except Exception:
        return None


def _extract_first_ppr(tf) -> etree._Element | None:
    """첫 번째 단락의 pPr(단락 서식) 요소 반환."""
    try:
        txBody = tf._txBody
        paras  = txBody.findall(qn("a:p"))
        if not paras:
            return None
        ppr = paras[0].find(qn("a:pPr"))
        return copy.deepcopy(ppr) if ppr is not None else None
    except Exception:
        return None


def _set_shape_text(shape, new_text: str) -> None:
    """
    도형의 텍스트를 교체하면서 원본 글꼴 서식(색상·크기·굵기)을 최대한 보존.
    도형의 위치·크기·배경 등 시각 요소는 건드리지 않음.
    """
    if not shape.has_text_frame:
        return

    tf     = shape.text_frame
    txBody = tf._txBody

    # 원본 서식 저장
    rpr = _extract_first_rpr(tf)
    ppr = _extract_first_ppr(tf)

    # 기존 단락 제거
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)

    lines = new_text.split("\n") if new_text else [""]
    for i, line in enumerate(lines):
        p_elem = etree.SubElement(txBody, qn("a:p"))

        # 단락 서식 (첫 단락에만 적용)
        if i == 0 and ppr is not None:
            p_elem.insert(0, copy.deepcopy(ppr))

        r_elem = etree.SubElement(p_elem, qn("a:r"))

        # 글꼴 서식 적용
        if rpr is not None:
            r_elem.insert(0, copy.deepcopy(rpr))

        t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line or " "

    tf.word_wrap = True


def _fill_slide(slide, title: str, body: str = "") -> None:
    """
    슬라이드의 텍스트 도형에 제목/본문을 채웁니다.
    - 텍스트 도형이 2개 이상: 첫 번째 = 제목, 두 번째 = 본문
    - 텍스트 도형이 1개: 제목 + 본문 합침
    - 텍스트 도형 없음: 변경 없음 (그래픽 전용 슬라이드)
    """
    shapes = _text_shapes_by_position(slide)
    if not shapes:
        return

    if len(shapes) == 1:
        combined = f"{title}\n{body}" if body else title
        _set_shape_text(shapes[0], combined)
    else:
        _set_shape_text(shapes[0], title)
        if body:
            _set_shape_text(shapes[1], body)
        # 3번째 이상 텍스트 도형은 비움
        for s in shapes[2:]:
            _set_shape_text(s, "")


# ══════════════════════════════════════════════════════════════════════════════
# 슬라이드 복제 (섹션 수 > 템플릿 슬라이드 수일 때)
# ══════════════════════════════════════════════════════════════════════════════

def _clone_slide(prs: Presentation, source_idx: int):
    """
    source_idx 슬라이드를 복제하여 프레젠테이션 끝에 추가.
    원본의 모든 도형(배경·그래픽·텍스트)을 그대로 복사.
    """
    source   = prs.slides[source_idx]
    layout   = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # 새 슬라이드의 기본 도형 트리 초기화
    sp_tree = new_slide.shapes._spTree
    required = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in list(sp_tree):
        if child.tag not in required:
            sp_tree.remove(child)

    # 원본 도형 복사 (required 태그 제외)
    src_sp_tree = source.shapes._spTree
    for child in src_sp_tree:
        if child.tag not in required:
            sp_tree.append(copy.deepcopy(child))

    return new_slide


# ══════════════════════════════════════════════════════════════════════════════
# 템플릿 기반 PPT 생성 (메인 로직)
# ══════════════════════════════════════════════════════════════════════════════

def _build_with_template(template_path: Path, analysis: dict) -> Presentation:
    prs = Presentation(str(template_path))
    n   = len(prs.slides)

    if n == 0:
        return _build_default(analysis)

    date_str    = analysis.get("date", datetime.now(timezone.utc).strftime("%Y-%m-%d"))
    provider    = analysis.get("provider", "-")
    attribution = analysis.get("model_attribution", provider)
    art_count   = analysis.get("article_count", 0)
    keywords    = ", ".join(analysis.get("keywords_found", [])[:6])
    report_type = analysis.get("type", "daily").upper()
    week_label  = analysis.get("week_label", "")
    subtitle    = week_label or date_str
    sections    = list(analysis.get("sections", {}).items())

    _safe_print(f"[PPT] 템플릿: {template_path.name} ({n}장) | 섹션: {len(sections)}개")

    # ── [슬라이드 1] 표지 ──────────────────────────────────────
    _fill_slide(
        prs.slides[0],
        title=(f"Automotive / AI Semiconductor\n"
               f"Strategy Brief  [{report_type}]"),
        body=(f"{subtitle}\n"
              f"수집 기사: {art_count}건  |  AI: {attribution}\n"
              f"키워드: {keywords}"),
    )

    # ── [슬라이드 2] 개요 (있을 경우) ─────────────────────────
    if n >= 2:
        sources = ", ".join(analysis.get("sources", [])[:8])
        _fill_slide(
            prs.slides[1],
            title="분석 개요 (Overview)",
            body=(f"■ 분석 일자: {subtitle}\n"
                  f"■ 수집 기사: {art_count}건\n"
                  f"■ 핵심 키워드: {keywords}\n"
                  f"■ 데이터 출처: {sources}\n"
                  f"■ AI 기여 모델: {attribution}"),
        )

    # ── [슬라이드 3~] 섹션 매핑 ───────────────────────────────
    # 섹션을 넣을 슬라이드 범위: index 2 ~ n-2 (마지막 슬라이드 보존)
    content_start = 2
    content_end   = max(n - 1, content_start)   # 마지막 슬라이드 제외
    available     = content_end - content_start  # 사용 가능한 슬라이드 수

    # 섹션 수가 가용 슬라이드 수보다 많으면 복제
    content_template_idx = content_start if content_start < n else 0
    while len(prs.slides) - 1 < content_start + len(sections):
        _clone_slide(prs, content_template_idx)
        _safe_print(f"[PPT] 슬라이드 복제 → 현재 {len(prs.slides)}장")

    # 섹션별로 슬라이드 채우기
    for i, (sec_title, sec_content) in enumerate(sections):
        slide_idx = content_start + i
        if slide_idx >= len(prs.slides):
            break
        body_text = _format_body(sec_content, max_lines=13)
        _fill_slide(prs.slides[slide_idx], title=sec_title, body=body_text)
        _safe_print(f"[PPT]  슬라이드 {slide_idx+1}: {sec_title[:40]}")

    # ── 마지막 슬라이드: 원본 유지 또는 출처 표시 ─────────────
    last_idx = len(prs.slides) - 1
    if last_idx > content_start + len(sections) - 1:
        sources_list = "\n".join(f"• {s}" for s in analysis.get("sources", [])[:12])
        _fill_slide(
            prs.slides[last_idx],
            title="데이터 출처 (Sources)",
            body=sources_list or "분석 완료",
        )

    return prs


# ══════════════════════════════════════════════════════════════════════════════
# 기본 다크 테마 (템플릿 없을 때 폴백)
# ══════════════════════════════════════════════════════════════════════════════

def _add_tb(slide, text, l, t, w, h, size=14, bold=False,
            color=None, align=PP_ALIGN.LEFT) -> None:
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE


def _set_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _build_default(analysis: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    date_str    = analysis.get("date", "")
    art_count   = analysis.get("article_count", 0)
    attribution = analysis.get("model_attribution", analysis.get("provider", "-"))
    report_type = analysis.get("type", "daily").upper()
    week_label  = analysis.get("week_label", "")
    subtitle    = week_label or date_str
    sections    = analysis.get("sections", {}) or {"분석결과": analysis.get("raw", "")}
    W = prs.slide_width

    # 표지
    s = prs.slides.add_slide(blank)
    _set_bg(s, C_BG)
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT; bar.line.fill.background()
    _add_tb(s, f"Automotive / AI Semiconductor  [{report_type}]",
            Inches(0.8), Inches(1.5), Inches(8.5), Inches(0.8),
            size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _add_tb(s, subtitle,
            Inches(0.8), Inches(2.5), Inches(8.5), Inches(0.5),
            size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_tb(s, f"수집 기사: {art_count}건  |  {attribution}",
            Inches(0.8), Inches(3.1), Inches(8.5), Inches(0.4),
            size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    slide_num = 2
    for title, content in sections.items():
        all_lines = []
        for raw in _clean_md(content).split("\n"):
            s2 = raw.strip().lstrip("-•·").strip()
            if s2:
                for wrapped in textwrap.wrap(s2, width=88):
                    all_lines.append("• " + wrapped)

        # 14줄씩 분할
        chunks = [all_lines[i:i+14] for i in range(0, len(all_lines), 14)] or [[]]
        for chunk_idx, chunk in enumerate(chunks):
            sl = prs.slides.add_slide(blank)
            _set_bg(sl, C_BG)
            bar2 = sl.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.07))
            bar2.fill.solid(); bar2.fill.fore_color.rgb = C_ACCENT; bar2.line.fill.background()
            disp = title if chunk_idx == 0 else f"{title} (계속 {chunk_idx+1})"
            _add_tb(sl, disp, Inches(0.3), Inches(0.1), Inches(9.2), Inches(0.6),
                    size=18, bold=True, color=C_ACCENT)
            box = sl.shapes.add_shape(1, Inches(0.3), Inches(0.85), Inches(9.2), Inches(6.2))
            box.fill.solid(); box.fill.fore_color.rgb = C_BODY_BG; box.line.fill.background()
            body_tb = sl.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(8.9), Inches(6.0))
            body_tf = body_tb.text_frame
            body_tf.word_wrap = True
            first = True
            for line in chunk:
                p = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
                first = False
                p.space_before = Pt(3)
                run = p.add_run()
                run.text = line
                run.font.size = Pt(12)
                run.font.color.rgb = C_LIGHT
            _add_tb(sl, str(slide_num), Inches(9.2), Inches(7.1), Inches(0.5), Inches(0.3),
                    size=10, color=RGBColor(0x44, 0x66, 0x88), align=PP_ALIGN.RIGHT)
            slide_num += 1

    return prs


# ══════════════════════════════════════════════════════════════════════════════
# Public API
# ══════════════════════════════════════════════════════════════════════════════

def generate_ppt(analysis: dict, output_dir: str = "/tmp") -> str:
    """
    분석 결과를 PPT로 저장.
    템플릿이 있으면 기존 슬라이드 디자인 유지, 없으면 기본 다크 테마.
    """
    date_str    = analysis.get("date", datetime.now(timezone.utc).strftime("%Y-%m-%d"))
    filename    = f"semiconductor_brief_{date_str}.pptx"
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    filepath    = str(output_path / filename)

    template = _find_template()
    if template:
        _safe_print(f"[PPT] 템플릿 사용: {template}")
        prs = _build_with_template(template, analysis)
    else:
        _safe_print("[PPT] 템플릿 없음 → 기본 다크 테마 사용")
        prs = _build_default(analysis)

    prs.save(filepath)
    _safe_print(f"[PPT] 저장 완료: {filepath}")
    return filepath
