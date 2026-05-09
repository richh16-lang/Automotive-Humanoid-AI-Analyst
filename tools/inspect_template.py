"""
PPT 템플릿 구조 분석 스크립트
실행: python tools/inspect_template.py templates/report_template.pptx

어떤 슬라이드 레이아웃과 플레이스홀더가 있는지 보여줍니다.
이 정보를 바탕으로 ppt_generator.py가 올바른 위치에 내용을 채워넣습니다.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

def inspect(path: str) -> None:
    prs = Presentation(path)
    print(f"\n══════════════════════════════════════════════")
    print(f"  템플릿 분석: {Path(path).name}")
    print(f"  슬라이드 크기: {prs.slide_width.inches:.1f}\" × {prs.slide_height.inches:.1f}\"")
    print(f"══════════════════════════════════════════════")

    # ── 슬라이드 레이아웃 목록 ────────────────────────────────
    print(f"\n[ 슬라이드 레이아웃 ({len(prs.slide_layouts)}개) ]")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n  레이아웃 [{i}]: '{layout.name}'")
        if layout.placeholders:
            for ph in layout.placeholders:
                idx  = ph.placeholder_format.idx
                ptype = str(ph.placeholder_format.type).split(".")[-1]
                try:
                    l = f"{ph.left.inches:.2f}\""
                    t = f"{ph.top.inches:.2f}\""
                    w = f"{ph.width.inches:.2f}\""
                    h = f"{ph.height.inches:.2f}\""
                    pos = f"위치: 좌{l} 상{t} 너비{w} 높이{h}"
                except Exception:
                    pos = ""
                print(f"    · 플레이스홀더 idx={idx:2d}  타입={ptype:15s}  {pos}")
        else:
            print("    (플레이스홀더 없음)")

    # ── 기존 슬라이드 목록 ────────────────────────────────────
    if prs.slides:
        print(f"\n[ 기존 슬라이드 ({len(prs.slides)}장) ]")
        for i, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name if slide.slide_layout else "-"
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()[:30]
                    if t:
                        texts.append(t)
            print(f"  슬라이드 [{i+1}] 레이아웃='{layout_name}'  텍스트: {texts[:2]}")

    # ── 권장 레이아웃 안내 ────────────────────────────────────
    print("\n══════════════════════════════════════════════")
    print("  권장 레이아웃 인덱스 설정 (.env에 추가)")
    print("  (직접 확인 후 맞는 번호로 설정하세요)")
    print()
    print("  TEMPLATE_PATH=templates/report_template.pptx")
    print("  TEMPLATE_TITLE_LAYOUT=0     # 표지 슬라이드 레이아웃 번호")
    print("  TEMPLATE_CONTENT_LAYOUT=1   # 내용 슬라이드 레이아웃 번호")
    print("  TEMPLATE_TITLE_PH=0         # 제목 플레이스홀더 idx")
    print("  TEMPLATE_BODY_PH=1          # 본문 플레이스홀더 idx")
    print("══════════════════════════════════════════════\n")

if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "templates/report_template.pptx"
    if not Path(path).exists():
        print(f"❌ 파일 없음: {path}")
        print("   템플릿 파일을 templates/ 폴더에 넣고 다시 실행하세요.")
        sys.exit(1)
    inspect(path)
